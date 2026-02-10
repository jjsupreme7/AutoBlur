"""PowerPoint parser using python-pptx."""

from typing import List

from pptx import Presentation

from .base import BaseParser, TextRegion, RedactionTarget


class PptxParser(BaseParser):
    def extract(self) -> List[TextRegion]:
        prs = Presentation(self.file_path)
        regions = []
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        text = run.text.strip()
                        if text:
                            regions.append(TextRegion(
                                text=run.text,
                                location={
                                    'slide_index': slide_idx,
                                    'shape_id': shape.shape_id,
                                    'paragraph_index': para_idx,
                                    'run_index': run_idx,
                                },
                                source_file=self.file_path,
                            ))
        return regions

    def redact(self, targets: List[RedactionTarget], output_path: str) -> None:
        prs = Presentation(self.file_path)
        for target in targets:
            loc = target.region.location
            slide = prs.slides[loc['slide_index']]
            for shape in slide.shapes:
                if shape.shape_id == loc['shape_id'] and shape.has_text_frame:
                    para = shape.text_frame.paragraphs[loc['paragraph_index']]
                    run = para.runs[loc['run_index']]
                    if target.region.text == run.text:
                        run.text = target.replacement
                    else:
                        run.text = run.text.replace(target.region.text, target.replacement)
        prs.save(output_path)
